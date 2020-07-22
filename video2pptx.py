# coding=utf-8
from tkinter import ttk
from tkinter import filedialog
from tkinter import messagebox
from tkinter import Tk
from tkinter import Label
from tkinter import Button
from tkinter import Entry
from tkinter import Canvas
from tkinter import Frame
from tkinter import PhotoImage
from tkinter import Toplevel
from tkinter import END
import time
import os.path
import threading
from queue import Queue
from pathlib import Path
import cv2
from skimage.metrics import structural_similarity as simi
import os
import pptx
from pptx.util import Inches
from pydub import AudioSegment
import inspect
import ctypes
import requests
import subprocess
import ffmpeg
import sys
import signal

image_arr = ['jpg', 'png', 'webp']
audio_arr = ['mp3', 'wav']
video_arr = ['rm', 'rmvb', 'flv', 'avi', 'mp4', 'mov']
audio_file = None
exception_state = 0
def _async_raise(tid, exctype):
    """raises the exception, performs cleanup if needed"""
    tid = ctypes.c_long(tid)
    if not inspect.isclass(exctype):
        exctype = type(exctype)
    res = ctypes.pythonapi.PyThreadState_SetAsyncExc(tid, ctypes.py_object(exctype))
    if res == 0:
        raise ValueError("invalid thread id")
    elif res != 1:
        # """if it returns a number greater than one, you're in trouble,
        # and you should call it again with exc=NULL to revert the effect"""
        ctypes.pythonapi.PyThreadState_SetAsyncExc(tid, None)
        raise SystemError("PyThreadState_SetAsyncExc failed")


def stop_thread(thread):
    if thread:
        _async_raise(thread.ident, SystemExit)


def get_video_info(video_path):
    global video_frames
    global video_width
    global video_height
    global video_fps
    global video_duration
    global video_size
    video_duration = 0
    video_size = 0
    try:
        probe = ffmpeg.probe(video_path)
        print(probe)
        if probe.get('format'):
            video_format = probe['format']
            if video_format.get('duration'):
                video_duration = video_format['duration']
            if video_format.get('size'):
                video_size = video_format['size']
        video_stream = next((stream for stream in probe['streams'] if stream['codec_type'] == 'video'), None)
        if video_stream is None:
            print('No video stream found', file=sys.stderr)
            return False
        video_width = int(video_stream['width'])
        video_height = int(video_stream['height'])
        display = video_stream['display_aspect_ratio']
        if video_stream.get('nb_frames'):
            video_frames = int(video_stream['nb_frames'])
        elif video_stream.get('avg_frame_rate'):
            video_frames = int(eval(video_stream['avg_frame_rate']) * int(float(video_duration)))
            video_stream['nb_frames'] = video_frames

        if video_stream.get('avg_frame_rate'):
            video_fps = eval(video_stream['avg_frame_rate'])

        # print('width: {}'.format(video_width))
        # print('height: {}'.format(video_height))
        # print('num_frames: {}'.format(video_frames))
        # print('display: ', display)
        # print('duration: ', video_duration)
        # print('size: ', video_size)
        return True
    except ffmpeg.Error as err:
        print(str(err.stderr, encoding='utf8'))
        return False


def post_record_aliyun(video_path, pptx_name, pptx_size, step_frames, sim_threshold, result):
    global video_frames
    global video_width
    global video_height
    global video_fps
    global video_duration
    global video_size
    url = 'http://106.15.227.254:5000/todos'
    payload = {
        "videoname": video_path,
        "videosize": video_size,
        "videoframes": video_frames,
        "duration": video_duration,
        "pptxname": pptx_name,
        "pptxsize": pptx_size,
        "step_frames": step_frames,
        "simi_threshold": sim_threshold,
        "result": result,
        "videofps": video_fps,
        "videoheight": video_height,
        "videowidth": video_width
    }
    response = requests.post(url=url, data=payload)
    # print(response.text)


def extract_full_audio(video_path, video_suffix, audio_format):
    # print(video_path, video_suffix, audio_format)
    if video_suffix is "flv":
        audio = AudioSegment.from_flv(video_path)
    else:
        audio = AudioSegment.from_file(video_path, video_suffix)
    audio.export(audio_file, format=audio_format)
    return audio


def on_closing():
    if messagebox.askokcancel("退出", "您确认要退出淘精助手吗？"):
        py.destroy()


def select_video():
    video_file = filedialog.askopenfilename(title='选择视频文件', filetypes=[('mp4', '*.mp4'), ('flv', '*.flv'),
                                                        ('avi', '*.avi'), ('mov', '*.mov'),
                                                        ('rmvb', '*.rmvb'), ('rm', '*.rm'),
                                                        ('所有支持文件', '*.mp4 *.flv *.avi *.mov *.rmvb *.rm')],
                                            initialdir=current_dir)
    entry_videopath['state'] = 'normal'
    entry_videopath.delete(0, END)
    entry_videopath.insert(0, video_file)
    entry_videopath['state'] = 'readonly'

    if entry_pptxname.get() is "":
        pptx_name = video_file
        suffix = pptx_name.split(".")
        # print(suffix)
        if len(suffix) > 1:
            pptx_name = pptx_name.replace(suffix[len(suffix)-1], "pptx")
            entry_pptxname.delete(0, END)
            entry_pptxname.insert(0, pptx_name)
    py.update()
    # print(video_file)


def set_pptxname():
    pptxname = filedialog.asksaveasfilename(title='设置PPTX文件目录和名称', filetypes=[('pptx', '*.pptx')],
                                            defaultextension=".pptx", initialdir=current_dir)
    suffix = pptxname.split(".")
    # print(suffix)
    if len(suffix) > 1:
        end_str = suffix[len(suffix) - 1]
        end_str = end_str.lower()
        if end_str != ".pptx":
            pptxname = pptxname.replace(suffix[len(suffix)-1], "pptx")
            info = "已经自动将PPTX文件" + pptxname + "后缀名由." + suffix[len(suffix)-1] + "改成了.pptx"
            messagebox.showinfo("提示", info)
    elif len(pptxname) > 0:
        pptxname = pptxname + ".pptx"

    entry_pptxname.delete(0, END)
    entry_pptxname.insert(0, pptxname)
    py.update()
    # print(pptxname)


def close_admire():
    top1.destroy()
    py.update()
    py.deiconify()


def on_sub_closing():
    time.sleep(0.2)


def show_wxpay_qrcode():
    global exception_state
    # print(exception_state, 1)
    if exception_state is 1:
        return

    global top1
    top1 = Toplevel()
    sub_height = 576
    sub_width = 424
    sub_pyx = (screen_width - sub_width) / 2
    sub_pyy = (screen_height - sub_height) / 2
    top1.iconbitmap('v2p.ico')
    top1.geometry('%dx%d+%d+%d' % (sub_width, sub_height, sub_pyx, sub_pyy))
    label_tips = Label(top1, text="恭喜--PPTX生成成功！淘精助手期待您的赞赏...", font="宋体 12", fg="red")
    label_tips.place(x=30, y=10, height=30)
    img_wxpay = PhotoImage(file='wxpay_qrcode.gif')
    label_wxpay = Label(top1, image=img_wxpay)
    label_wxpay.place(x=30, y=50, height=459, width=358)
    button_skip = Button(top1, text="跳过", width=10, state='normal', command=close_admire)
    button_skip.place(x=180, y=530, height=30)
    top1.minsize(sub_width, sub_height)
    top1.maxsize(sub_width, sub_height)
    top1.protocol("WM_DELETE_WINDOW", on_sub_closing)
    py.withdraw()
    top1.mainloop()


def extract_frames_split_audio(queue, video_path, step_frames, sim_threshold, image_format,
                               audio_format, output_path, audio_file):
    for i in range(400):
        queue.put(10+i)
    shell_cmd = 'ffmpeg -v 0 -i ' + video_path + ' -f ' + audio_format + ' ' + audio_file + ' -r ' +\
                str(float(video_fps)/int(step_frames)) + ' -ss 0.5 -f image2 ' + output_path + '/%d.' + image_format
    print(shell_cmd)
    # global ffmpid
    subprocess.call(shell_cmd, shell=True)
    # print(ffmpid, ffmpid.pid)
    pre = None
    preffn = None
    apreidx = 0
    acuridx = 0
    jpgs = [fn for fn in os.listdir(output_path) if fn.endswith(image_format)]
    # print(jpgs)
    for i in range(180):
        queue.put(410 + i)
    # print(output_path)
    for fn in sorted(jpgs, key=lambda item: int(item[:item.rindex('.')])):
        ffn = output_path + "/" + fn
        if (pre is None) and (os.path.exists(ffn)):
            preffn = ffn
            fnn = fn.split(".")
            if fnn[0].isdigit():
                apreidx = int(fnn[0])
            pre = cv2.imread(ffn)
            continue
        if os.path.exists(ffn):
            cur = cv2.imread(ffn)

        fnn = fn.split(".")
        if fnn[0].isdigit():
            acuridx = int(fnn[0])

        if simi(pre, cur, multichannel=True) < float(sim_threshold):
            start = (apreidx-1) * int(step_frames) / float(video_fps)
            duration = (acuridx - apreidx) * int(step_frames) / float(video_fps)
            afn = preffn.replace(image_format, audio_format)
            # print(afn)

            shell_cmd = 'ffmpeg -v 0 -i ' + audio_file + ' -ss ' + str(start) + ' -t ' + str(duration) + ' -f ' +\
                        audio_format + ' ' + afn
            print(shell_cmd)
            subprocess.call(shell_cmd, shell=True)
            preffn = ffn
            apreidx = acuridx
            pre = cur
        else:
            os.remove(ffn)

    if apreidx < acuridx:
        start = apreidx * int(step_frames) / float(video_fps)
        duration = float(video_duration) - start
        afn = preffn.replace(image_format, audio_format)
        # print(afn)
        shell_cmd = 'ffmpeg -v 0 -i ' + audio_file + ' -ss ' + str(start) + ' -t ' + str(duration) + ' -f ' + \
                    audio_format + ' ' + afn
        print(shell_cmd)
        subprocess.call(shell_cmd, shell=True)
    queue.put(600)


def extract_images_cut_audios(queue, vc, audio, step_frames, sim_threshold, image_format, audio_format, output_path):
    pre = None
    preidx = 0
    times = int(int(video_frames)/int(step_frames))
    for i in range(times):
        curidx = i * int(step_frames)
        if curidx > int(video_frames):
            break
        vc.set(cv2.CAP_PROP_POS_FRAMES, curidx)  # 设置要获取的帧号
        ret, frame = vc.read()
        if frame is None:
            break
        if pre is None:
            cv2.imwrite(os.path.join(output_path, '{}.{}'.format(curidx, image_format)), frame)
            pre = frame
            continue
        if simi(pre, frame, multichannel=True) < float(sim_threshold):
            cv2.imwrite(os.path.join(output_path, '{}.{}'.format(curidx, image_format)), frame)
            if curidx:
                start = (1000.0 * preidx) / int(video_fps)
                end = (1000.0 * (curidx - 1)) / int(video_fps)
                audiofile = audio[start:end]
                audiofile.export(os.path.join(output_path, '{}.{}'.format(preidx, audio_format)), format=audio_format)
            pre = frame
            preidx = curidx
        queue.put(10 + int((i * 590) / times))

    if preidx < video_frames:
        start = (1000.0 * preidx) / video_fps
        end = (1000.0 * video_frames) / video_fps
        audiofile = audio[start:end]
        audiofile.export(os.path.join(output_path, '{}.{}'.format(preidx, audio_format)), format=audio_format)
    queue.put(600)


def write_pptx(queue, image_format, audio_format, pptx_name, output_path):
    ppts = pptx.Presentation()
    ppts.slide_height = 6858000
    ppts.slide_width = 12192000

    jpgs = [fn for fn in os.listdir(output_path) if fn.endswith(image_format)]
    cover = True
    idx = 0
    for fn in sorted(jpgs, key=lambda item: int(item[:item.rindex('.')])):
        if cover:
            slide = ppts.slides.add_slide(ppts.slide_layouts[0])
            cover = False
        else:
            slide = ppts.slides.add_slide(ppts.slide_layouts[6])
        ffn = output_path + "/" + fn
        afn = ffn.replace(image_format, audio_format)
        print(afn)
        if os.path.exists(ffn):
            slide.shapes.add_picture(ffn, Inches(0), Inches(0), width=ppts.slide_width)
        if os.path.exists(afn):
            slide.shapes.add_movie(afn, Inches(0), Inches(0), Inches(1), Inches(1), poster_frame_image=None,
                               mime_type='video/mp4')
        queue.put(610 + int(idx*80/len(jpgs)))
        idx = idx + 1
    ppts.save(pptx_name)
    queue.put(699)


def del_files(path_data):
    path_exist = Path(path_data)
    if path_exist.is_dir():
        for root, dirs, files in os.walk(path_data, topdown=False):
            for name in files:
                os.remove(os.path.join(root, name))
            for name in dirs:
                os.rmdir(os.path.join(root, name))
        #os.remove(path_data)


def remove_dir_files(output_path):
    if os.path.exists(output_path):
        del_files(output_path)
    global audio_file
    if audio_file and os.path.exists(audio_file):
        os.remove(audio_file)


def video2pptx(process_queue, video_path, pptx_name, sim_threshold, output_path,
               step_frames, image_format='jpg', audio_format='mp3'):
    # print(video_path, pptx_name, sim_threshold, output_path, step_frames, image_format, audio_format)
    global exception_state
    video_file = Path(video_path)
    if video_file.is_file() is False:
        msg = "视频文件" + video_path + "不存在，请检查！"
        process_queue.put(msg)
        return
    process_queue.put(1)

    suffix = video_path.split(".")
    # print(suffix)
    if len(suffix) <= 1:
        msg = "视频文件" + video_path + "不存在后缀名，请检查！"
        process_queue.put(msg)
        return
    process_queue.put(2)

    video_suffix = suffix[len(suffix) - 1]
    global audio_file
    audio_file = video_path.replace(video_suffix, audio_format)
    # print(audio_file)
    if os.path.exists(audio_file):
        os.remove(audio_file)

    pptx_dir = Path(os.path.dirname(pptx_name))
    if pptx_dir.is_dir() is False:
        msg = "PPTX文件夹" + pptx_dir + "不存在，请检查！"
        process_queue.put(msg)
        return
    if os.path.exists(pptx_name):
        os.remove(pptx_name)

    process_queue.put(3)
    if (float(sim_threshold) > 0.95) or (float(sim_threshold) < 0.85):
        msg = "图像相似度阈值有效取值范围为[85%, 95%]，当前取值为" + sim_threshold + "%，请修改!"
        process_queue.put(msg)
        return
    process_queue.put(4)

    out_dir = Path(output_path)
    if out_dir.is_dir() is False:
        # 去除首位空格
        path = output_path.strip()
        # 去除尾部 \ 符号
        path = path.rstrip("\\")
        os.makedirs(path)
    process_queue.put(5)

    if (int(step_frames)%100>0) and (int(step_frames) < 100 or int(step_frames) > 3000):
        msg = "步进帧数的取值为(100,200,300,400,500,600,700,800,900,1000,2000,30000)，当前取值为"+step_frames + "请修改!"
        process_queue.put(msg)
        return
    process_queue.put(6)

    if image_format not in image_arr:
        msg = "选择的图片格式" + image_format + "不在淘精助手支持的图片格式" + image_arr + "中，请修改!"
        process_queue.put(msg)
        return
    process_queue.put(7)

    if audio_format not in audio_arr:
        msg = "选择的音频格式" + audio_format + "不在淘精助手支持的音频格式" + audio_arr + "中，请修改!"
        process_queue.put(msg)
        return
    process_queue.put(8)

    if exception_state:
        return

    ret = get_video_info(video_path)
    if ret is False:
        msg = "抱歉，淘精助手不支持视频文件" + video_path + "的音视频编码，无法完成转换!"
        process_queue.put(msg)
        return
    process_queue.put(9)
    '''
    audio = extract_full_audio(video_path, video_suffix, audio_format)
    process_queue.put(10)

    extract_images_cut_audios(process_queue, vc, audio, step_frames, sim_threshold, image_format,
                              audio_format, output_path)
    vc.release()
    '''
    if exception_state:
        return
    extract_frames_split_audio(process_queue, video_path, step_frames, sim_threshold, image_format,
                               audio_format, output_path, audio_file)
    process_queue.put(610)
    if exception_state:
        return

    write_pptx(process_queue, image_format, audio_format, pptx_name, output_path)
    if exception_state:
        return
    pptx_size = os.path.getsize(pptx_name)
    process_queue.put(690)
    remove_dir_files(output_path)
    if exception_state:
        return
    post_record_aliyun(video_path, pptx_name, pptx_size, step_frames, sim_threshold, "convert successfully")
    if exception_state:
        return
    process_queue.put(700)


def stop_except_convert():
    fill_line = canvas.create_rectangle(1.5, 1.5, 0, 30, width=0, fill="white")
    canvas.coords(fill_line, (0, 0, 701, 60))
    button_quit['state'] = 'normal'
    button_stop['state'] = 'disabled'
    button_convert['state'] = 'normal'
    button_selectvideo['state'] = 'normal'
    button_pptxname['state'] = 'normal'
    entry_pptxname['state'] = 'normal'
    combo_stepframe['state'] = 'readonly'
    combo_simithreshold['state'] = 'readonly'
    py.update()
    output_path = current_dir + "\\output\\"
    remove_dir_files(output_path)


def progress_update(queue):
    # 填充进度条
    fill_line = canvas.create_rectangle(1.5, 1.5, 0, 30, width=0, fill="green")
    preidx = 0
    global process_state
    process_state = 1
    while process_state:
        if queue.empty():
            idx = preidx + 0.1
        else:
            msg = str(queue.get_nowait())
            if msg.isdigit():
                idx = int(msg)
            else:
                global exception_state
                exception_state = 1
                video_path = entry_videopath.get()
                if video_path is not "":
                    pptx_name = entry_pptxname.get()
                    step_frames = combo_stepframe.get()
                    sim_threshold = int(combo_simithreshold.get()) / 100.0
                    post_record_aliyun(video_path, pptx_name, 0, step_frames, sim_threshold, msg)
                messagebox.showinfo("提示", msg)
                stop_except_convert()
                break
            preidx = idx
        if canvas:
            canvas.coords(fill_line, (0, 0, idx, 60))
        if idx >= 699:
            if canvas:
                canvas.coords(fill_line, (0, 0, 700, 60))
            break
        py.update()
        time.sleep(1)  # 控制进度条流动的速度
    py.update()


def start_convert():
    global exception_state
    exception_state = 0
    button_selectvideo['state'] = 'disabled'
    button_pptxname['state'] = 'disabled'
    button_convert['state'] = 'disabled'
    button_stop['state'] = 'normal'
    button_quit['state'] = 'disabled'
    entry_pptxname['state'] = 'readonly'
    combo_stepframe['state'] = 'disabled'
    combo_simithreshold['state'] = 'disabled'
    # 清空进度条
    fill_line = canvas.create_rectangle(1.5, 1.5, 0, 30, width=0, fill="white")
    canvas.coords(fill_line, (0, 0, 701, 60))
    py.update()

    video_path = entry_videopath.get()
    pptx_name = entry_pptxname.get()
    step_frames = combo_stepframe.get()
    simi_threshold = int(combo_simithreshold.get())/100.0
    output_path = current_dir + "\\output\\"
    process_queue = Queue(maxsize=1000)
    # 先创造线程
    global process_thread
    process_thread = threading.Thread(target=video2pptx, args=(process_queue, video_path, pptx_name,
                                                               simi_threshold, output_path, step_frames,
                                                               "jpg", "mp3"))
    # progress_thread = threading.Thread(target=progress_update, args=(process_queue,))

    start_time = time.time()
    # 启动线程
    # process_thread.setDaemon(True)
    process_thread.start()
    progress_update(process_queue)
    print("last time: {} s".format(time.time() - start_time))

    button_quit['state'] = 'normal'
    button_stop['state'] = 'disabled'
    button_convert['state'] = 'normal'
    button_selectvideo['state'] = 'normal'
    button_pptxname['state'] = 'normal'
    entry_pptxname['state'] = 'normal'
    combo_stepframe['state'] = 'readonly'
    combo_simithreshold['state'] = 'readonly'
    py.update()
    if exception_state:
        return
    show_wxpay_qrcode()


def stop_convert():
    # global ffmpid
    global process_state
    global exception_state
    global process_thread
    exception_state = 1
    process_state = 0

    os.system("taskkill /IM ffmpeg.exe /F")
    # stop_thread(process_thread)
    # process_thread.join()

    time.sleep(1)
    # 清空进度条
    fill_line = canvas.create_rectangle(1.5, 1.5, 0, 30, width=0, fill="white")
    canvas.coords(fill_line, (0, 0, 701, 60))
    button_quit['state'] = 'normal'
    button_stop['state'] = 'disabled'
    button_convert['state'] = 'normal'
    button_selectvideo['state'] = 'normal'
    button_pptxname['state'] = 'normal'
    entry_pptxname['state'] = 'normal'
    combo_stepframe['state'] = 'readonly'
    combo_simithreshold['state'] = 'readonly'
    py.update()
    output_path = current_dir + "\\output\\"

    remove_dir_files(output_path)


def quit_convert():
    # show_wxpay_qrcode()
    py.destroy()


if __name__ == "__main__":
    py = Tk()
    py.withdraw()
    window_height = 576
    window_width = 1024

    current_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = current_dir + "\\output\\"
    if os.path.exists(output_path):
        del_files(output_path)
    py.title("淘精助手--视频转PPTX V1.0.0")  # 设置标题
    py.iconbitmap('v2p.ico')  # 设置图标
    py.minsize(window_width, window_height)
    py.maxsize(window_width, window_height)

    screen_height = py.winfo_screenheight()
    screen_width = py.winfo_screenwidth()

    pyx = (screen_width - window_width) / 2
    pyy = (screen_height - window_height) / 2
    # print(pyx)
    # print(pyy)
    # print(sys.getdefaultencoding())

    py.geometry('%dx%d-%d+%d' % (window_width, window_height, pyx, pyy))
    py.protocol("WM_DELETE_WINDOW", on_closing)

    # 视频文件路径
    Label(py, text="视频文件路径").place(x=30, y=30, height=30)
    entry_videopath = Entry(py)
    entry_videopath.place(x=140, y=30, height=30, width=700)
    entry_videopath['state'] = 'readonly'

    button_selectvideo = Button(py, text="选择视频", width=10, state='normal', command=select_video)
    button_selectvideo.place(x=880, y=30)

    # pptx文件名称
    Label(py, text="PPTX文件名称").place(x=30, y=80, height=30)
    entry_pptxname = Entry(py)
    entry_pptxname.place(x=140, y=80, height=30, width=700)

    button_pptxname = Button(py, text="设置文件", width=10, state='normal', command=set_pptxname)
    button_pptxname.place(x=880, y=80)

    # 步进帧数
    Label(py, text="步进帧数").place(x=30, y=130, height=30)
    combo_stepframe = ttk.Combobox(py, state="readonly")
    combo_stepframe.place(x=140, y=130, width=70, height=30)
    combo_stepframe['values'] = (100, 200, 300, 400, 500, 600, 700, 800, 900, 1000, 2000, 3000,10000)
    combo_stepframe.current(8)
    Label(py, text="帧").place(x=220, y=130, height=30)

    # 图像相似度阈值
    Label(py, text="图像相似度阈值").place(x=30, y=180, height=30)
    combo_simithreshold = ttk.Combobox(py, state="readonly")
    combo_simithreshold.place(x=140, y=180, width=70, height=30)
    combo_simithreshold['values'] = (95, 94, 93, 92, 91, 90, 89, 88, 87, 86, 85)
    combo_simithreshold.current(0)
    Label(py, text="%").place(x=220, y=180, height=30)

    # 转换进度条
    Label(py, text="转换进度条").place(x=30, y=280, height=30)
    canvas = Canvas(py, bg="white")
    canvas.place(x=140, y=280, width=700, height=30)

    # 命令按钮
    button_convert = Button(py, text="开始转换", width=10, state='normal', command=start_convert)
    button_convert.place(x=300, y=230)
    button_stop = Button(py, text="停止转换", width=10, state='disabled', command=stop_convert)
    button_stop.place(x=480, y=230)
    button_quit = Button(py, text="退出", width=10, state='normal', command=quit_convert)
    button_quit.place(x=650, y=230)

    Frame(height=230, width=964, bd=1, background='black').place(x=30, y=330)
    Frame(height=228, width=962, bd=1).place(x=31, y=331)
    usetext = "使用说明：\n1、点击‘选择视频’按钮，选择要转换的视频文件，目前支持的视频格式有rm、rmvb、flv、mp4、" + \
        "avi、mov；\n2、在弹出的“选择视频文件”对话框中选择视频文件；\n3、选择完视频文件后，PPTX文件名称会自动填充" + \
        "，如果要修改PPTX文件名称，可以直接在编辑框中进行修改，也可以通过点击‘设置文件’按钮修改文件名称和路径；" + \
        "\n4、在下拉列表框中，选择‘步进帧数’，步进帧数默认值是900，是用于调整每次移动的帧数，用户可以根据实际情况" + \
        "选择需要的步进帧数；\n5、在下拉列表框中，选择‘图像相似度阈值’，图像相似度阈值默认为95" + \
        "，一般不需要调整；\n6、点击‘开始转换’按钮，开始将视频文件转存成PPTX文件，整个转换过程取决于电脑性能、" + \
        "视频大小和分辨率，一般情况下1个小时的视频可以在15分钟之内完成；\n7、如果转换过程中，要停止转换" + \
        "，请点击‘停止转换’按钮；\n8、如果要退出淘精助手，请点击‘退出’按钮。\n\n备注：使用中有问题或相关合作，" + \
        "请发邮件到1606863058@qq.com 或 直接加QQ 1606863058"
    Label(py, text=usetext, anchor="w", justify="left").place(x=36, y=336, height=218, width=952)

    global process_state
    process_state = 0
    py.deiconify()
    py.mainloop()

